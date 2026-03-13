#!/usr/bin/env python3
"""
PPT 渲染引擎 — 读取 JSON 配置文件，生成专业 .pptx 文件。

运行环境：Docker 沙箱，Python 3.11
依赖：python-pptx（自动安装）
中文字体：WenQuanYi Micro Hei（文泉驿微米黑）
"""

import os
import sys
import subprocess
import json
import argparse
import re
import traceback

# ---------------------------------------------------------------------------
# 自动安装 python-pptx
# ---------------------------------------------------------------------------
try:
    from pptx import Presentation
except ImportError:
    skill_dir = os.environ.get('SKILL_DIR', '')
    pkg_dir = os.path.join(skill_dir, 'packages')
    if os.path.isdir(pkg_dir):
        subprocess.check_call(
            [sys.executable, '-m', 'pip', 'install',
             '--no-index', '--find-links', pkg_dir, 'python-pptx', '-q'],
        )
    else:
        subprocess.check_call(
            [sys.executable, '-m', 'pip', 'install', 'python-pptx', '-q'],
        )
    from pptx import Presentation

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Chart imports
try:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    CategoryChartData = None
    XL_CHART_TYPE = None

# ---------------------------------------------------------------------------
# 常量
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 通用边距
MARGIN_LEFT = Inches(0.8)
MARGIN_RIGHT = Inches(0.8)
MARGIN_TOP = Inches(0.4)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT

# 标题区高度
TITLE_AREA_HEIGHT = Inches(1.0)
# 标题线偏移
TITLE_LINE_TOP = MARGIN_TOP + TITLE_AREA_HEIGHT + Inches(0.05)
# 内容区起始
CONTENT_TOP = TITLE_LINE_TOP + Inches(0.25)
CONTENT_AVAILABLE_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - Inches(0.8)

# 页脚区域
FOOTER_TOP = SLIDE_HEIGHT - Inches(0.45)
FOOTER_HEIGHT = Inches(0.3)

# ---------------------------------------------------------------------------
# 字体检测：优先微软雅黑，降级到 Noto Sans CJK SC
# ---------------------------------------------------------------------------
def _detect_cjk_font():
    """检测系统可用的中文字体，优先微软雅黑，降级到 Noto Sans CJK SC"""
    try:
        result = subprocess.run(['fc-list', ':lang=zh'], capture_output=True, text=True, timeout=5)
        if '微软雅黑' in result.stdout or 'Microsoft YaHei' in result.stdout:
            return '微软雅黑'
        if 'Noto Sans CJK SC' in result.stdout:
            return 'Noto Sans CJK SC'
    except Exception:
        pass
    return '微软雅黑'  # fallback

CJK_FONT = _detect_cjk_font()

# ---------------------------------------------------------------------------
# 主题定义
# ---------------------------------------------------------------------------
THEMES = {
    'sgcc': {
        'name': '国网绿',
        'primary': RGBColor(0, 101, 80),
        'secondary': RGBColor(196, 163, 90),
        'accent': RGBColor(0, 61, 165),
        'bg': RGBColor(255, 255, 255),
        'text': RGBColor(51, 51, 51),
        'light_text': RGBColor(255, 255, 255),
        'title_font': CJK_FONT,
        'body_font': CJK_FONT,
        'divider_bg': RGBColor(0, 101, 80),
        'footer_text': '内部资料 注意保密',
        'neutral_line': RGBColor(200, 200, 200),
        'stripe_even': RGBColor(255, 255, 255),
        'stripe_odd': RGBColor(242, 242, 242),
    },
    'professional': {
        'name': '商务蓝',
        'primary': RGBColor(0, 61, 165),
        'secondary': RGBColor(245, 245, 245),
        'accent': RGBColor(220, 53, 69),
        'bg': RGBColor(255, 255, 255),
        'text': RGBColor(51, 51, 51),
        'light_text': RGBColor(255, 255, 255),
        'title_font': CJK_FONT,
        'body_font': CJK_FONT,
        'divider_bg': RGBColor(0, 61, 165),
        'footer_text': '',
        'neutral_line': RGBColor(200, 200, 200),
        'stripe_even': RGBColor(255, 255, 255),
        'stripe_odd': RGBColor(242, 242, 242),
    },
    'minimal': {
        'name': '简约',
        'primary': RGBColor(51, 51, 51),
        'secondary': RGBColor(245, 245, 245),
        'accent': RGBColor(0, 122, 204),
        'bg': RGBColor(255, 255, 255),
        'text': RGBColor(68, 68, 68),
        'light_text': RGBColor(255, 255, 255),
        'title_font': CJK_FONT,
        'body_font': CJK_FONT,
        'divider_bg': RGBColor(51, 51, 51),
        'footer_text': '',
        'neutral_line': RGBColor(200, 200, 200),
        'stripe_even': RGBColor(255, 255, 255),
        'stripe_odd': RGBColor(242, 242, 242),
    },
    'tech': {
        'name': '科技蓝',
        'primary': RGBColor(0, 120, 212),
        'secondary': RGBColor(30, 30, 46),
        'accent': RGBColor(0, 200, 83),
        'bg': RGBColor(255, 255, 255),
        'text': RGBColor(51, 51, 51),
        'light_text': RGBColor(255, 255, 255),
        'title_font': CJK_FONT,
        'body_font': CJK_FONT,
        'divider_bg': RGBColor(0, 120, 212),
        'footer_text': '',
        'neutral_line': RGBColor(200, 200, 200),
        'stripe_even': RGBColor(255, 255, 255),
        'stripe_odd': RGBColor(242, 242, 242),
    },
}

# 图表类型映射
CHART_TYPE_MAP = {}
if XL_CHART_TYPE is not None:
    CHART_TYPE_MAP = {
        'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
    }


# =========================================================================
# 工具函数
# =========================================================================

def set_chinese_font(run, font_name=CJK_FONT):
    """同时设置拉丁字体和东亚字体（中文需要通过 XML 设置 a:ea 标签）"""
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font_name)


def add_styled_textbox(slide, left, top, width, height, text, theme,
                       font_size=Pt(18), bold=False, color=None,
                       alignment=None, font_key='body_font'):
    """创建带样式的文本框，减少重复代码。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color or theme['text']
    set_chinese_font(run, theme.get(font_key, theme['body_font']))
    return txBox, tf, p, run


def render_bullet_list(tf, bullets, theme, max_count=8, start_new=True,
                       bullet_size=Pt(12), text_size=Pt(18),
                       sub_bullet_size=Pt(10), sub_text_size=Pt(15)):
    """渲染 bullet 列表到 TextFrame，支持子 bullet。"""
    if max_count:
        bullets = bullets[:max_count]

    for idx, bullet_text in enumerate(bullets):
        sub_bullets = []
        if isinstance(bullet_text, dict):
            sub_bullets = bullet_text.get('sub_bullets', [])
            bullet_text = bullet_text.get('text', str(bullet_text))

        if idx == 0 and not start_new:
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        else:
            p = tf.add_paragraph() if (idx > 0 or start_new) else tf.paragraphs[0]

        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.line_spacing = 1.5
        p.level = 0

        # Bullet 符号
        bullet_run = p.add_run()
        bullet_run.text = '\u25cf  '
        bullet_run.font.size = bullet_size
        bullet_run.font.color.rgb = theme['primary']
        set_chinese_font(bullet_run, theme['body_font'])

        # 正文
        text_run = p.add_run()
        text_run.text = str(bullet_text)
        text_run.font.size = text_size
        text_run.font.color.rgb = theme['text']
        set_chinese_font(text_run, theme['body_font'])

        # 子 bullet
        for sub in sub_bullets:
            sp = tf.add_paragraph()
            sp.space_before = Pt(2)
            sp.space_after = Pt(2)
            sp.line_spacing = 1.4
            sp.level = 1

            sp_run_indent = sp.add_run()
            sp_run_indent.text = '      \u25cb  '
            sp_run_indent.font.size = sub_bullet_size
            sp_run_indent.font.color.rgb = theme['accent']
            set_chinese_font(sp_run_indent, theme['body_font'])

            sp_run = sp.add_run()
            sp_run.text = str(sub)
            sp_run.font.size = sub_text_size
            sp_run.font.color.rgb = theme['text']
            set_chinese_font(sp_run, theme['body_font'])


def add_divider_line(slide, left, top, width, height, color):
    """添加装饰线/分隔线。"""
    line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = color
    line_shape.line.fill.background()
    line_shape.shadow.inherit = False
    return line_shape


def add_text_run(tf, text, font_name=CJK_FONT, size=Pt(18),
                 bold=False, italic=False, color=None, new_paragraph=False):
    """向 TextFrame 中追加一段带格式的文本。

    Args:
        tf: TextFrame 对象
        text: 文本内容
        font_name: 字体名称
        size: 字号
        bold: 是否粗体
        italic: 是否斜体
        color: RGBColor 颜色，None 表示不设置
        new_paragraph: True 则新建段落，否则在当前段落追加 run
    Returns:
        (paragraph, run) 元组
    """
    if new_paragraph and len(tf.paragraphs) > 0 and tf.paragraphs[0].text != '':
        p = tf.add_paragraph()
    elif new_paragraph and len(tf.paragraphs) > 0 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.paragraphs[-1] if tf.paragraphs else tf.add_paragraph()

    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    set_chinese_font(run, font_name)
    return p, run


def add_shaped_rect(slide, left, top, width, height, fill_color):
    """添加一个填充矩形，无边框。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # 无边框
    shape.shadow.inherit = False
    return shape


def add_title_bar(slide, title_text, theme):
    """在幻灯片顶部添加标准标题区（标题 + 装饰线）。

    Returns:
        标题 shape 的底部 y 坐标，方便后续布局
    """
    # 标题文本
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, TITLE_AREA_HEIGHT,
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = theme['primary']
    set_chinese_font(run, theme['title_font'])
    p.alignment = PP_ALIGN.LEFT

    # 装饰线
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, TITLE_LINE_TOP,
        CONTENT_WIDTH, Pt(2),
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = theme['primary']
    line_shape.line.fill.background()
    line_shape.shadow.inherit = False

    return CONTENT_TOP


def sanitize_filename(name):
    """将文件名中的特殊字符替换为下划线。"""
    return re.sub(r'[\\/:*?"<>|]', '_', name).strip()


def generate_theme_colors(theme, count):
    """根据主题生成 count 个图表配色。

    使用 primary、accent、secondary 以及插值色来扩展。
    """
    base_colors = [theme['primary'], theme['accent']]
    # 添加一些衍生颜色
    extra_colors = [
        RGBColor(75, 172, 198),   # 青色
        RGBColor(255, 165, 0),    # 橙色
        RGBColor(128, 0, 128),    # 紫色
        RGBColor(34, 139, 34),    # 森林绿
        RGBColor(220, 20, 60),    # 猩红
        RGBColor(100, 149, 237),  # 矢车菊蓝
    ]
    colors = base_colors + extra_colors
    # 循环填充
    result = []
    for i in range(count):
        result.append(colors[i % len(colors)])
    return result


# =========================================================================
# 幻灯片渲染函数
# =========================================================================

def render_title_slide(prs, slide_data, theme):
    """渲染封面页。

    布局：顶部 60% 区域填充主题色，放置主副标题；
    底部 40% 白色区域放置日期和作者信息。
    """
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    top_height = int(SLIDE_HEIGHT * 0.6)
    bottom_height = SLIDE_HEIGHT - top_height

    # 顶部主题色背景矩形
    add_shaped_rect(slide, 0, 0, SLIDE_WIDTH, top_height, theme['primary'])

    # 主标题
    title_text = slide_data.get('title', '演示文稿')
    title_box = slide.shapes.add_textbox(
        Inches(1.5), int(top_height * 0.28),
        SLIDE_WIDTH - Inches(3.0), Inches(1.5),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = theme['light_text']
    set_chinese_font(run, theme['title_font'])

    # 副标题
    subtitle_text = slide_data.get('subtitle', '')
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(
            Inches(1.5), int(top_height * 0.28) + Inches(1.6),
            SLIDE_WIDTH - Inches(3.0), Inches(1.0),
        )
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.size = Pt(20)
        run2.font.color.rgb = theme['light_text']
        set_chinese_font(run2, theme['body_font'])

    # 底部白色区域 — 日期和作者
    bottom_top = top_height + int(bottom_height * 0.3)
    info_parts = []
    if slide_data.get('author'):
        info_parts.append(slide_data['author'])
    if slide_data.get('date'):
        info_parts.append(slide_data['date'])
    if slide_data.get('department'):
        info_parts.append(slide_data['department'])

    if info_parts:
        info_box = slide.shapes.add_textbox(
            Inches(1.5), bottom_top,
            SLIDE_WIDTH - Inches(3.0), Inches(1.0),
        )
        tf3 = info_box.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = '  |  '.join(info_parts)
        run3.font.size = Pt(16)
        run3.font.color.rgb = theme['text']
        set_chinese_font(run3, theme['body_font'])

    # 底部装饰线（金色 / accent）
    deco_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.0), top_height - Pt(4),
        SLIDE_WIDTH - Inches(6.0), Pt(4),
    )
    deco_line.fill.solid()
    deco_line.fill.fore_color.rgb = theme['secondary']
    deco_line.line.fill.background()
    deco_line.shadow.inherit = False

    return slide


def render_agenda_slide(prs, slide_data, theme):
    """渲染目录 / 议程页。

    显示编号列表，支持高亮当前章节。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_text = slide_data.get('title', '目录')
    add_title_bar(slide, title_text, theme)

    # 议程项目
    items = slide_data.get('items', [])
    highlight_index = slide_data.get('highlight', -1)

    if not items:
        return slide

    item_height = min(Inches(0.8), int(CONTENT_AVAILABLE_HEIGHT / max(len(items), 1)))
    start_top = CONTENT_TOP + Inches(0.2)

    for idx, item in enumerate(items):
        item_text = item if isinstance(item, str) else item.get('text', str(item))
        is_highlighted = (idx == highlight_index)

        item_left = MARGIN_LEFT + Inches(0.3)
        item_top = start_top + idx * item_height
        item_width = CONTENT_WIDTH - Inches(0.6)

        # 高亮背景
        if is_highlighted:
            bg_shape = add_shaped_rect(
                slide, item_left - Inches(0.15), item_top,
                item_width + Inches(0.3), item_height - Inches(0.08),
                theme['primary'],
            )
            # 确保背景在底层（z-order 较低即可，python-pptx 按添加顺序叠放）

        # 编号圆形
        circle_size = Inches(0.45)
        circle_left = item_left
        circle_top = item_top + (item_height - circle_size) // 2 - Inches(0.04)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, circle_left, circle_top,
            circle_size, circle_size,
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = theme['light_text'] if is_highlighted else theme['primary']
        circle.line.fill.background()
        circle.shadow.inherit = False

        # 编号数字
        circle_tf = circle.text_frame
        circle_tf.word_wrap = False
        circle_p = circle_tf.paragraphs[0]
        circle_p.alignment = PP_ALIGN.CENTER
        circle_run = circle_p.add_run()
        circle_run.text = str(idx + 1)
        circle_run.font.size = Pt(16)
        circle_run.font.bold = True
        circle_run.font.color.rgb = theme['primary'] if is_highlighted else theme['light_text']
        set_chinese_font(circle_run, theme['title_font'])
        circle_tf.paragraphs[0].space_before = Pt(0)
        circle_tf.paragraphs[0].space_after = Pt(0)

        # 项目文本
        text_left = circle_left + circle_size + Inches(0.3)
        text_box = slide.shapes.add_textbox(
            text_left, item_top,
            item_width - circle_size - Inches(0.6), item_height,
        )
        ttf = text_box.text_frame
        ttf.word_wrap = True
        ttf.paragraphs[0].space_before = Pt(0)
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        tr = tp.add_run()
        tr.text = item_text
        tr.font.size = Pt(20)
        tr.font.bold = is_highlighted
        tr.font.color.rgb = theme['light_text'] if is_highlighted else theme['text']
        set_chinese_font(tr, theme['body_font'])

        # 垂直居中
        text_box.text_frame.paragraphs[0].space_before = Pt(8)

    return slide


def render_section_divider_slide(prs, slide_data, theme):
    """渲染章节分隔页。

    整页填充 divider_bg 色，标题和副标题居中。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 全页背景
    add_shaped_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, theme['divider_bg'])

    # 左侧装饰竖线
    deco = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), int(SLIDE_HEIGHT * 0.3),
        Pt(6), int(SLIDE_HEIGHT * 0.4),
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = theme['secondary']
    deco.line.fill.background()
    deco.shadow.inherit = False

    # 章节标题
    section_title = slide_data.get('title', '')
    title_box = slide.shapes.add_textbox(
        Inches(2.2), int(SLIDE_HEIGHT * 0.32),
        SLIDE_WIDTH - Inches(4.0), Inches(1.8),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = section_title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = theme['light_text']
    set_chinese_font(run, theme['title_font'])

    # 副标题
    subtitle = slide_data.get('subtitle', '')
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(2.2), int(SLIDE_HEIGHT * 0.32) + Inches(1.6),
            SLIDE_WIDTH - Inches(4.0), Inches(1.0),
        )
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(18)
        run2.font.color.rgb = theme['light_text']
        set_chinese_font(run2, theme['body_font'])

    # 章节编号（如有）
    section_number = slide_data.get('section_number', '')
    if section_number:
        num_box = slide.shapes.add_textbox(
            SLIDE_WIDTH - Inches(2.5), int(SLIDE_HEIGHT * 0.35),
            Inches(1.5), Inches(1.5),
        )
        ntf = num_box.text_frame
        np_ = ntf.paragraphs[0]
        np_.alignment = PP_ALIGN.RIGHT
        nr = np_.add_run()
        nr.text = str(section_number)
        nr.font.size = Pt(72)
        nr.font.bold = True
        nr.font.color.rgb = theme['secondary']
        set_chinese_font(nr, theme['title_font'])

    return slide


def render_content_slide(prs, slide_data, theme):
    """渲染正文内容页。

    标题 + 装饰线 + bullet 列表（最多显示内容区能容纳的条目数）。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_text = slide_data.get('title', '')
    add_title_bar(slide, title_text, theme)

    bullets = slide_data.get('bullets', [])
    if not bullets:
        bullets = slide_data.get('content', [])
    if isinstance(bullets, str):
        bullets = [bullets]

    # Bullet 区域
    bullet_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(0.2), CONTENT_TOP + Inches(0.1),
        CONTENT_WIDTH - Inches(0.4), CONTENT_AVAILABLE_HEIGHT,
    )
    tf = bullet_box.text_frame
    tf.word_wrap = True

    render_bullet_list(tf, bullets, theme, max_count=8, start_new=False)

    return slide


def render_two_column_slide(prs, slide_data, theme):
    """渲染双栏内容页。

    标题区 + 左右两栏（各 45% 宽度，10% 间距），
    每栏有小标题和 bullet 列表。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_text = slide_data.get('title', '')
    add_title_bar(slide, title_text, theme)

    # 栏宽计算
    gap = Inches(0.6)
    col_width = int((CONTENT_WIDTH - gap) / 2)
    left_col_left = MARGIN_LEFT
    right_col_left = MARGIN_LEFT + col_width + gap

    # 支持两种数据格式：
    # 格式 A（嵌套）: left: { title, bullets }, right: { title, bullets }
    # 格式 B（扁平）: left_title, left_bullets, right_title, right_bullets
    left_data = slide_data.get('left', {})
    right_data = slide_data.get('right', {})
    if not left_data and not right_data:
        left_data = {
            'title': slide_data.get('left_title', ''),
            'bullets': slide_data.get('left_bullets', []),
        }
        right_data = {
            'title': slide_data.get('right_title', ''),
            'bullets': slide_data.get('right_bullets', []),
        }

    for col_idx, (col_data, col_left) in enumerate([
        (left_data, left_col_left),
        (right_data, right_col_left),
    ]):
        col_title = col_data.get('heading', col_data.get('title', ''))
        col_bullets = col_data.get('bullets', [])
        if isinstance(col_bullets, str):
            col_bullets = [col_bullets]

        current_top = CONTENT_TOP + Inches(0.1)

        # 栏标题
        if col_title:
            add_styled_textbox(
                slide, col_left, current_top, col_width, Inches(0.6),
                col_title, theme,
                font_size=Pt(20), bold=True, color=theme['primary'],
                font_key='title_font',
            )
            current_top += Inches(0.7)

            # 栏标题下划线
            add_divider_line(
                slide, col_left, current_top - Inches(0.15),
                col_width, Pt(1.5), theme['primary'],
            )

        # 栏内容
        if col_bullets:
            bullets_box = slide.shapes.add_textbox(
                col_left + Inches(0.1), current_top + Inches(0.05),
                col_width - Inches(0.2),
                CONTENT_AVAILABLE_HEIGHT - Inches(0.8),
            )
            btf = bullets_box.text_frame
            btf.word_wrap = True

            render_bullet_list(
                btf, col_bullets, theme, max_count=6, start_new=False,
                bullet_size=Pt(10), text_size=Pt(16),
            )

    # 中间分隔线（可选）
    if slide_data.get('divider', True):
        divider_x = MARGIN_LEFT + col_width + int(gap / 2)
        add_divider_line(
            slide, divider_x, CONTENT_TOP,
            Pt(1), CONTENT_AVAILABLE_HEIGHT - Inches(0.3),
            theme.get('neutral_line', RGBColor(200, 200, 200)),
        )

    return slide


def render_chart_slide(prs, slide_data, theme):
    """渲染图表页。

    支持 bar / line / pie / doughnut 图表类型，
    使用 python-pptx 原生图表 API。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_text = slide_data.get('title', '')
    add_title_bar(slide, title_text, theme)

    chart_type_str = slide_data.get('chart_type', 'bar').lower()
    if chart_type_str not in CHART_TYPE_MAP:
        print(f'[WARN] 不支持的图表类型 "{chart_type_str}"，使用 bar 替代',
              file=sys.stderr)
        chart_type_str = 'bar'

    chart_type = CHART_TYPE_MAP.get(chart_type_str, CHART_TYPE_MAP.get('bar'))

    # 支持 data.categories / data.series（SKILL.md 规范）和顶层字段（兼容）
    data_obj = slide_data.get('data', {})
    categories = data_obj.get('categories', slide_data.get('categories', []))
    series_list = data_obj.get('series', slide_data.get('series', []))

    if not categories or not series_list:
        # 无数据时显示提示
        placeholder = slide.shapes.add_textbox(
            Inches(3), CONTENT_TOP + Inches(1.5),
            SLIDE_WIDTH - Inches(6), Inches(1.0),
        )
        ptf = placeholder.text_frame
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = '[ 图表数据为空 ]'
        pr.font.size = Pt(20)
        pr.font.color.rgb = RGBColor(180, 180, 180)
        set_chinese_font(pr, theme['body_font'])
        return slide

    chart_data = CategoryChartData()
    # 数据类型校验
    if not isinstance(categories, list):
        categories = []
    if not isinstance(series_list, list):
        series_list = []
    chart_data.categories = categories
    for s in series_list:
        if not isinstance(s, dict):
            continue
        s_name = s.get('name', '系列')
        s_values = s.get('values', [])
        # 确保 values 都是数字
        s_values = [v if isinstance(v, (int, float)) else 0 for v in s_values]
        chart_data.add_series(s_name, s_values)

    # 图表位置和大小
    chart_left = MARGIN_LEFT + Inches(0.5)
    chart_top = CONTENT_TOP + Inches(0.2)
    chart_width = CONTENT_WIDTH - Inches(1.0)
    chart_height = CONTENT_AVAILABLE_HEIGHT - Inches(0.3)

    chart_frame = slide.shapes.add_chart(
        chart_type, chart_left, chart_top, chart_width, chart_height,
        chart_data,
    )

    chart = chart_frame.chart
    chart.has_legend = len(series_list) > 1

    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)

    # 应用主题颜色到系列
    colors = generate_theme_colors(theme, len(series_list))
    for idx, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[idx]

        # 折线图设置线条颜色
        if chart_type_str == 'line':
            series.format.line.color.rgb = colors[idx]
            series.format.line.width = Pt(2.5)
            series.smooth = False

    # 饼图 / 环形图配色（按 point）
    if chart_type_str in ('pie', 'doughnut') and len(chart.series) > 0:
        point_colors = generate_theme_colors(theme, len(categories))
        series_obj = chart.series[0]
        for pidx in range(len(categories)):
            try:
                point = series_obj.points[pidx]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = point_colors[pidx]
            except (IndexError, AttributeError):
                pass

    # 图表样式微调
    try:
        if hasattr(chart, 'category_axis'):
            chart.category_axis.tick_labels.font.size = Pt(11)
        if hasattr(chart, 'value_axis'):
            chart.value_axis.tick_labels.font.size = Pt(11)
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(220, 220, 220)
    except Exception:
        pass  # 饼图没有 category_axis / value_axis

    return slide


def render_table_slide(prs, slide_data, theme):
    """渲染表格页。

    表头行使用主题色背景 + 白色文字，数据行斑马纹。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_text = slide_data.get('title', '')
    add_title_bar(slide, title_text, theme)

    headers = slide_data.get('headers', [])
    rows_data = slide_data.get('rows', [])

    if not headers and not rows_data:
        return slide

    num_rows = len(rows_data) + 1  # +1 for header
    num_cols = len(headers) if headers else (len(rows_data[0]) if rows_data else 1)

    # 表格尺寸和位置（居中）
    table_width = min(CONTENT_WIDTH - Inches(0.4), Inches(11.5))
    table_height = min(
        CONTENT_AVAILABLE_HEIGHT - Inches(0.2),
        Inches(0.5) * num_rows + Inches(0.2),
    )
    table_left = MARGIN_LEFT + int((CONTENT_WIDTH - table_width) / 2)
    table_top = CONTENT_TOP + Inches(0.15)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, table_left, table_top, table_width, table_height,
    )
    table = table_shape.table

    # 列宽均分
    col_width = int(table_width / num_cols)
    for col_idx in range(num_cols):
        table.columns[col_idx].width = col_width

    # 行高
    header_row_height = Inches(0.55)
    data_row_height = Inches(0.45)
    table.rows[0].height = header_row_height
    for r_idx in range(1, num_rows):
        table.rows[r_idx].height = data_row_height

    # 填充表头
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = ''
        # 清除默认段落，手动设置
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(header)
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = theme['light_text']
        set_chinese_font(run, theme['title_font'])

        # 表头背景色
        _set_cell_bg(cell, theme['primary'])

        # 垂直居中
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        # 内边距
        _set_cell_margins(cell, Inches(0.1), Inches(0.05), Inches(0.1), Inches(0.05))

    # 补齐/截断行数据，确保与 headers 列数一致
    for r_idx in range(len(rows_data)):
        row = rows_data[r_idx]
        if not isinstance(row, list):
            rows_data[r_idx] = [str(row)] + [''] * (num_cols - 1)
        elif len(row) < num_cols:
            rows_data[r_idx] = list(row) + [''] * (num_cols - len(row))
        elif len(row) > num_cols:
            rows_data[r_idx] = row[:num_cols]

    # 填充数据行
    stripe_colors = [theme.get('stripe_even', RGBColor(255, 255, 255)),
                     theme.get('stripe_odd', RGBColor(242, 242, 242))]
    for r_idx, row_values in enumerate(rows_data):
        stripe_color = stripe_colors[(r_idx) % 2]
        for col_idx in range(num_cols):
            cell = table.cell(r_idx + 1, col_idx)
            cell.text = ''
            value = row_values[col_idx] if col_idx < len(row_values) else ''

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(value)
            run.font.size = Pt(13)
            run.font.color.rgb = theme['text']
            set_chinese_font(run, theme['body_font'])

            _set_cell_bg(cell, stripe_color)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            _set_cell_margins(cell, Inches(0.1), Inches(0.04), Inches(0.1), Inches(0.04))

    return slide


def _set_cell_bg(cell, color):
    """设置表格单元格背景色。"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # 移除已有的 solidFill
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {
        'val': str(color),
    })
    solidFill.append(srgbClr)
    tcPr.insert(0, solidFill)


def _set_cell_margins(cell, left, top, right, bottom):
    """设置表格单元格内边距。"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set('marL', str(int(left)))
    tcPr.set('marT', str(int(top)))
    tcPr.set('marR', str(int(right)))
    tcPr.set('marB', str(int(bottom)))


def render_image_text_slide(prs, slide_data, theme):
    """渲染图文混排页。

    左 60% 放图片，右 40% 放文字。
    图片不存在时显示灰色占位矩形。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_text = slide_data.get('title', '')
    add_title_bar(slide, title_text, theme)

    # 布局参数
    image_width = int(CONTENT_WIDTH * 0.58)
    text_width = int(CONTENT_WIDTH * 0.38)
    image_left = MARGIN_LEFT
    text_left = MARGIN_LEFT + image_width + int(CONTENT_WIDTH * 0.04)
    area_top = CONTENT_TOP + Inches(0.15)
    area_height = CONTENT_AVAILABLE_HEIGHT - Inches(0.2)

    # 图片区域
    image_path = slide_data.get('image', '')
    # 相对于 cwd (workspace) 解析路径
    if image_path and not os.path.isabs(image_path):
        workspace = os.environ.get('WORKSPACE_PATH', os.getcwd())
        image_path = os.path.join(workspace, image_path)

    # 路径安全检查：确保图片在 workspace 内
    if image_path:
        workspace = os.environ.get('WORKSPACE_PATH', os.getcwd())
        real_image = os.path.realpath(image_path)
        real_workspace = os.path.realpath(workspace)
        if not real_image.startswith(real_workspace + os.sep) and real_image != real_workspace:
            print(f'[WARN] 图片路径逃逸 workspace: {image_path}', file=sys.stderr)
            image_path = ''  # 阻止加载

    if image_path and os.path.isfile(image_path):
        try:
            pic = slide.shapes.add_picture(
                image_path, image_left, area_top, image_width, area_height,
            )
            # 保持纵横比
            img_ratio = pic.image.size[0] / pic.image.size[1]
            box_ratio = image_width / area_height
            if img_ratio > box_ratio:
                # 图片更宽，以宽度为准
                new_width = image_width
                new_height = int(image_width / img_ratio)
                pic.width = new_width
                pic.height = new_height
                pic.top = area_top + int((area_height - new_height) / 2)
            else:
                # 图片更高，以高度为准
                new_height = area_height
                new_width = int(area_height * img_ratio)
                pic.width = new_width
                pic.height = new_height
                pic.left = image_left + int((image_width - new_width) / 2)
        except Exception as e:
            print(f'[WARN] 加载图片失败: {e}', file=sys.stderr)
            _add_image_placeholder(slide, image_left, area_top,
                                   image_width, area_height, theme)
    else:
        _add_image_placeholder(slide, image_left, area_top,
                               image_width, area_height, theme,
                               text=f'图片未找到: {slide_data.get("image", "N/A")}')

    # 文字区域
    text_content = slide_data.get('text', '')
    text_bullets = slide_data.get('bullets', [])

    text_box = slide.shapes.add_textbox(
        text_left, area_top, text_width, area_height,
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    if text_content:
        p = tf.paragraphs[0]
        p.line_spacing = 1.6
        run = p.add_run()
        run.text = str(text_content)
        run.font.size = Pt(16)
        run.font.color.rgb = theme['text']
        set_chinese_font(run, theme['body_font'])

    for bidx, bt in enumerate(text_bullets):
        if isinstance(bt, dict):
            bt = bt.get('text', str(bt))
        np = tf.add_paragraph() if (text_content or bidx > 0) else tf.paragraphs[0]
        np.space_before = Pt(4)
        np.line_spacing = 1.5

        dot = np.add_run()
        dot.text = '\u25cf  '
        dot.font.size = Pt(10)
        dot.font.color.rgb = theme['primary']
        set_chinese_font(dot, theme['body_font'])

        tr = np.add_run()
        tr.text = str(bt)
        tr.font.size = Pt(15)
        tr.font.color.rgb = theme['text']
        set_chinese_font(tr, theme['body_font'])

    return slide


def _add_image_placeholder(slide, left, top, width, height, theme, text='图片占位'):
    """添加灰色占位矩形代替缺失的图片。"""
    placeholder = add_shaped_rect(
        slide, left, top, width, height, RGBColor(230, 230, 230),
    )
    # 占位文字
    placeholder.text_frame.word_wrap = True
    p = placeholder.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(160, 160, 160)
    set_chinese_font(run, theme['body_font'])
    placeholder.text_frame.paragraphs[0].space_before = Pt(int(height / Pt(1) / 3))


def render_quote_slide(prs, slide_data, theme):
    """渲染引用页。

    居中布局，带装饰性大引号、引文和来源。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 淡色背景
    add_shaped_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, theme['bg'])

    quote_text = slide_data.get('quote', slide_data.get('text', ''))
    source_text = slide_data.get('source', slide_data.get('author', ''))

    # 装饰性左引号
    quote_mark_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.8),
        Inches(2.0), Inches(2.0),
    )
    qm_tf = quote_mark_box.text_frame
    qm_p = qm_tf.paragraphs[0]
    qm_run = qm_p.add_run()
    qm_run.text = '\u201c'  # 左双引号 "
    qm_run.font.size = Pt(120)
    qm_run.font.bold = True
    qm_run.font.color.rgb = theme['primary']
    set_chinese_font(qm_run, theme['title_font'])

    # 引文
    quote_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(2.5),
        SLIDE_WIDTH - Inches(5.0), Inches(3.0),
    )
    qtf = quote_box.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.alignment = PP_ALIGN.CENTER
    qp.line_spacing = 1.8
    qr = qp.add_run()
    qr.text = str(quote_text)
    qr.font.size = Pt(24)
    qr.font.italic = True
    qr.font.color.rgb = theme['text']
    set_chinese_font(qr, theme['body_font'])

    # 来源
    if source_text:
        # 装饰线
        src_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(SLIDE_WIDTH / 2) - Inches(1.0), Inches(5.5),
            Inches(2.0), Pt(2),
        )
        src_line.fill.solid()
        src_line.fill.fore_color.rgb = theme['primary']
        src_line.line.fill.background()
        src_line.shadow.inherit = False

        src_box = slide.shapes.add_textbox(
            Inches(2.5), Inches(5.7),
            SLIDE_WIDTH - Inches(5.0), Inches(0.6),
        )
        stf = src_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.RIGHT
        sr = sp.add_run()
        sr.text = f'\u2014\u2014 {source_text}'  # —— 来源
        sr.font.size = Pt(16)
        sr.font.color.rgb = theme['text']
        set_chinese_font(sr, theme['body_font'])

    return slide


def render_thank_you_slide(prs, slide_data, theme):
    """渲染结束页 / 致谢页。

    整页填充主题色背景，居中大字 + 可选联系方式。
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 全页背景
    add_shaped_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, theme['primary'])

    # 主文字
    main_text = slide_data.get('title', slide_data.get('text', '谢谢'))
    main_box = slide.shapes.add_textbox(
        Inches(2.0), Inches(2.2),
        SLIDE_WIDTH - Inches(4.0), Inches(2.0),
    )
    tf = main_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = main_text
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = theme['light_text']
    set_chinese_font(run, theme['title_font'])

    # 副标题（如有）
    subtitle = slide_data.get('subtitle', '')
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(2.0), Inches(4.0),
            SLIDE_WIDTH - Inches(4.0), Inches(1.0),
        )
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(20)
        sr.font.color.rgb = theme['light_text']
        set_chinese_font(sr, theme['body_font'])

    # 联系方式
    contact = slide_data.get('contact', '')
    if contact:
        # 分隔线
        deco = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(SLIDE_WIDTH / 2) - Inches(1.5), Inches(5.0),
            Inches(3.0), Pt(2),
        )
        deco.fill.solid()
        deco.fill.fore_color.rgb = theme['secondary']
        deco.line.fill.background()
        deco.shadow.inherit = False

        contact_box = slide.shapes.add_textbox(
            Inches(2.0), Inches(5.3),
            SLIDE_WIDTH - Inches(4.0), Inches(1.0),
        )
        ctf = contact_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(contact)
        cr.font.size = Pt(18)
        cr.font.color.rgb = theme['light_text']
        set_chinese_font(cr, theme['body_font'])

    return slide


# =========================================================================
# 布局分发器
# =========================================================================

LAYOUT_RENDERERS = {
    'title': render_title_slide,
    'cover': render_title_slide,
    'agenda': render_agenda_slide,
    'toc': render_agenda_slide,
    'section': render_section_divider_slide,
    'section_divider': render_section_divider_slide,
    'divider': render_section_divider_slide,
    'content': render_content_slide,
    'bullets': render_content_slide,
    'two_column': render_two_column_slide,
    'two_columns': render_two_column_slide,
    'comparison': render_two_column_slide,
    'chart': render_chart_slide,
    'table': render_table_slide,
    'image_text': render_image_text_slide,
    'image': render_image_text_slide,
    'quote': render_quote_slide,
    'thank_you': render_thank_you_slide,
    'thanks': render_thank_you_slide,
    'end': render_thank_you_slide,
}


# =========================================================================
# 页脚 & 页码
# =========================================================================

def add_slide_numbers(prs, theme, skip_first=True, skip_last=True):
    """为所有幻灯片添加页脚和页码。

    跳过封面页（第一页）和结束页（最后一页）。
    """
    total = len(prs.slides)
    footer_text = theme.get('footer_text', '')
    footer_color = RGBColor(160, 160, 160)

    for idx, slide in enumerate(prs.slides):
        # 跳过封面和结束页
        if skip_first and idx == 0:
            continue
        if skip_last and idx == total - 1:
            continue

        # 左侧页脚文字
        if footer_text:
            left_box = slide.shapes.add_textbox(
                MARGIN_LEFT, FOOTER_TOP,
                Inches(4.0), FOOTER_HEIGHT,
            )
            ltf = left_box.text_frame
            lp = ltf.paragraphs[0]
            lp.alignment = PP_ALIGN.LEFT
            lr = lp.add_run()
            lr.text = footer_text
            lr.font.size = Pt(10)
            lr.font.color.rgb = footer_color
            set_chinese_font(lr, theme['body_font'])

        # 右侧页码
        right_box = slide.shapes.add_textbox(
            SLIDE_WIDTH - MARGIN_RIGHT - Inches(2.0), FOOTER_TOP,
            Inches(2.0), FOOTER_HEIGHT,
        )
        rtf = right_box.text_frame
        rp = rtf.paragraphs[0]
        rp.alignment = PP_ALIGN.RIGHT
        rr = rp.add_run()
        rr.text = f'{idx + 1} / {total}'
        rr.font.size = Pt(10)
        rr.font.color.rgb = footer_color
        set_chinese_font(rr, theme['body_font'])


# =========================================================================
# 配置加载 & 校验
# =========================================================================

def parse_args():
    """解析命令行参数，获取 JSON 配置文件路径。"""
    parser = argparse.ArgumentParser(
        description='PPT 渲染引擎 — 从 JSON 配置生成 .pptx 文件',
    )
    parser.add_argument(
        'config',
        help='JSON 配置文件路径',
    )
    parser.add_argument(
        '-o', '--output',
        help='输出文件路径（覆盖默认的 OUTPUTS_PATH）',
        default=None,
    )
    parser.add_argument(
        '-t', '--theme',
        help='主题名称（覆盖配置文件中的 theme 字段）',
        default=None,
    )
    return parser.parse_args()


def load_config(config_path):
    """加载并校验 JSON 配置文件。

    Returns:
        解析后的 dict
    Raises:
        SystemExit: JSON 格式错误或缺少必要字段
    """
    if not os.path.isfile(config_path):
        print(json.dumps({
            'success': False,
            'error': f'配置文件不存在: {config_path}',
        }, ensure_ascii=False))
        sys.exit(1)

    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            config = json.load(f)
    except json.JSONDecodeError as e:
        print(json.dumps({
            'success': False,
            'error': f'JSON 格式错误: {e}',
        }, ensure_ascii=False))
        sys.exit(1)

    # 基本校验
    if 'slides' not in config or not isinstance(config['slides'], list):
        print(json.dumps({
            'success': False,
            'error': '配置文件缺少 slides 数组',
        }, ensure_ascii=False))
        sys.exit(1)

    if len(config['slides']) == 0:
        print(json.dumps({
            'success': False,
            'error': 'slides 数组为空',
        }, ensure_ascii=False))
        sys.exit(1)

    return config


def get_theme(config, override_theme=None):
    """获取主题配色方案。

    Args:
        config: 配置 dict
        override_theme: 命令行指定的主题（优先级最高）
    Returns:
        主题 dict
    """
    metadata = config.get('metadata', {})
    theme_name = override_theme or metadata.get('theme', config.get('theme', 'professional'))
    theme_name = theme_name.lower().strip()

    if theme_name not in THEMES:
        print(f'[WARN] 未知主题 "{theme_name}"，使用 professional 替代',
              file=sys.stderr)
        theme_name = 'professional'

    return THEMES[theme_name]


# =========================================================================
# 主流程
# =========================================================================

def create_presentation(config, theme):
    """创建 Presentation 对象，遍历 slides 配置逐个渲染。

    Args:
        config: 配置 dict
        theme: 主题 dict
    Returns:
        pptx.Presentation 对象
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slides_config = config.get('slides', [])
    rendered_count = 0
    error_count = 0

    for idx, slide_data in enumerate(slides_config):
        layout = slide_data.get('layout', 'content').lower().strip()

        renderer = LAYOUT_RENDERERS.get(layout)
        if renderer is None:
            print(f'[WARN] 第 {idx + 1} 页布局 "{layout}" 不支持，已跳过',
                  file=sys.stderr)
            continue

        try:
            renderer(prs, slide_data, theme)
            rendered_count += 1
        except Exception as e:
            error_count += 1
            print(f'[WARN] 第 {idx + 1} 页渲染失败: {e}', file=sys.stderr)
            traceback.print_exc(file=sys.stderr)
            # 尝试添加一张错误提示页
            try:
                error_slide_layout = prs.slide_layouts[6]
                error_slide = prs.slides.add_slide(error_slide_layout)
                err_box = error_slide.shapes.add_textbox(
                    Inches(2), Inches(3), Inches(9), Inches(1.5),
                )
                etf = err_box.text_frame
                ep = etf.paragraphs[0]
                ep.alignment = PP_ALIGN.CENTER
                er = ep.add_run()
                er.text = f'渲染错误: {e}'
                er.font.size = Pt(18)
                er.font.color.rgb = RGBColor(220, 53, 69)
                set_chinese_font(er, theme['body_font'])
                rendered_count += 1
            except Exception:
                pass

    return prs, rendered_count, error_count


def save_presentation(prs, config, output_override=None):
    """保存 PPTX 文件到 OUTPUTS_PATH。

    Returns:
        保存的文件路径
    """
    # 确定输出目录
    if output_override:
        output_path = output_override
    else:
        outputs_dir = os.environ.get('OUTPUTS_PATH', os.path.join(os.getcwd(), 'outputs'))
        os.makedirs(outputs_dir, exist_ok=True)

        # 文件名
        metadata = config.get('metadata', {})
        title = metadata.get('title', config.get('title', '演示文稿'))
        filename = sanitize_filename(title) + '.pptx'
        output_path = os.path.join(outputs_dir, filename)

    # 确保目录存在
    out_dir = os.path.dirname(output_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)

    prs.save(output_path)
    return output_path


def main():
    """主入口：解析参数 → 加载配置 → 渲染 → 添加页码 → 保存。"""
    args = parse_args()

    # 加载配置
    config = load_config(args.config)

    # 获取主题
    theme = get_theme(config, args.theme)

    # 用 metadata.footer_text 覆盖主题默认页脚
    metadata = config.get('metadata', {})
    if metadata.get('footer_text'):
        theme = dict(theme)  # 浅拷贝避免修改全局 THEMES
        theme['footer_text'] = metadata['footer_text']

    # 创建演示文稿
    prs, slide_count, error_count = create_presentation(config, theme)

    if slide_count == 0:
        print(json.dumps({
            'success': False,
            'error': '没有成功渲染任何幻灯片',
        }, ensure_ascii=False))
        sys.exit(1)

    # 添加页脚和页码
    add_slide_numbers(prs, theme)

    # 保存
    output_path = save_presentation(prs, config, args.output)

    # 计算相对路径用于输出
    workspace = os.environ.get('WORKSPACE_PATH', os.getcwd())
    try:
        rel_path = os.path.relpath(output_path, workspace)
    except ValueError:
        rel_path = output_path

    # 输出结果
    result = {
        'success': True,
        'file': rel_path,
        'absolute_path': output_path,
        'slides': slide_count,
        'errors': error_count,
        'theme': theme['name'],
    }
    print(json.dumps(result, ensure_ascii=False))


if __name__ == '__main__':
    main()
